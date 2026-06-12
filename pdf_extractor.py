import PyPDF2
from docx import Document
from pptx import Presentation
import io

class DocumentExtractor:
    @staticmethod
    def extract_from_pdf(file_obj):
        """Extract text from PDF"""
        try:
            pdf_reader = PyPDF2.PdfReader(file_obj)
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page.extract_text()
            return text
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"
    
    @staticmethod
    def extract_from_docx(file_obj):
        """Extract text from DOCX"""
        try:
            doc = Document(file_obj)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text += row_text + "\n"
            
            return text
        except Exception as e:
            return f"Error extracting DOCX: {str(e)}"
    
    @staticmethod
    def extract_from_pptx(file_obj):
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_obj)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"Error extracting PPTX: {str(e)}"
    
    @staticmethod
    def extract_from_txt(file_obj):
        """Extract text from TXT"""
        try:
            return file_obj.read().decode('utf-8')
        except Exception as e:
            return f"Error extracting TXT: {str(e)}"
    
    @staticmethod
    def extract_text(file_obj, file_type):
        """Main extraction method"""
        if file_type.lower() == 'pdf':
            return DocumentExtractor.extract_from_pdf(file_obj)
        elif file_type.lower() == 'docx':
            return DocumentExtractor.extract_from_docx(file_obj)
        elif file_type.lower() == 'pptx':
            return DocumentExtractor.extract_from_pptx(file_obj)
        elif file_type.lower() == 'txt':
            return DocumentExtractor.extract_from_txt(file_obj)
        else:
            return "Unsupported file type"
